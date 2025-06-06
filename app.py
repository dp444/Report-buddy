import os
import PyPDF2
import pandas as pd
import docx
import fitz  # PyMuPDF
from flask import Flask, request, jsonify, send_from_directory
from flask_cors import CORS
from docx import Document as DocxDocument
from werkzeug.utils import secure_filename
from PIL import Image
from io import BytesIO
from google import genai
from google.genai import types
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import markdown2
from bs4 import BeautifulSoup

# ---------- Gemini API Setup ----------
client = genai.Client(api_key="")#Add google gemini api token here

# ---------- Flask Setup ----------
app = Flask(__name__)
CORS(app)

UPLOAD_FOLDER = "uploaded_files"
TEMPLATE_FOLDER = "template_reports"
OUTPUT_FOLDER = "summaries"

os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(TEMPLATE_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)

# ---------- File Parsers ----------
def extract_text_from_pdf(file_path):
    text = ""
    with open(file_path, "rb") as f:
        reader = PyPDF2.PdfReader(f)
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    return text

def extract_text_from_docx(file_path):
    doc = docx.Document(file_path)
    return "\n".join([p.text for p in doc.paragraphs])

def extract_text_from_excel(file_path):
    data = pd.read_excel(file_path)
    return data.to_string(index=False)

def extract_headings_from_docx(file_path):
    doc = docx.Document(file_path)
    return [p.text.strip() for p in doc.paragraphs if p.style.name.startswith("Heading") and p.text.strip()]

# ---------- Image Extractors ----------
def extract_images_from_pdf(pdf_path, output_folder):
    doc = fitz.open(pdf_path)
    image_paths = []
    for page_num in range(len(doc)):
        page = doc[page_num]
        images = page.get_images(full=True)
        for i, img in enumerate(images):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            image_ext = base_image["ext"]
            image_filename = f"pdf_page{page_num+1}_img{i+1}.{image_ext}"
            image_path = os.path.join(output_folder, image_filename)
            with open(image_path, "wb") as f:
                f.write(image_bytes)
            image_paths.append(image_path)
    return image_paths

def extract_images_from_docx(docx_path, output_folder):
    doc = DocxDocument(docx_path)
    image_paths = []
    rels = doc.part._rels
    for rel in rels:
        rel = rels[rel]
        if "image" in rel.target_ref:
            img_data = rel.target_part.blob
            image_filename = os.path.basename(rel.target_ref)
            image_path = os.path.join(output_folder, image_filename)
            with open(image_path, "wb") as f:
                f.write(img_data)
            image_paths.append(image_path)
    return image_paths

# ---------- Gemini Generators ----------
def generate_report_text_with_gemini(text, headings=None):
    prompt = f"""
Create a detailed ppt in bullet/section format based on the following content:

{text}

If the data includes tables or numeric values:
- Include bullet points and descriptions
- Mention recommended charts like bar/line/pie (in text)

{f"Use this structure:\n" + ''.join(f"- {h}\n" for h in headings) if headings else ""}
"""
    response = client.models.generate_content(model='gemini-2.0-flash', contents=[prompt])
    return response.text

def generate_image_with_gemini(prompt):
    try:
        response = client.models.generate_content(
            model="gemini-2.0-flash-exp-image-generation",
            contents=prompt,
            config=types.GenerateContentConfig(response_modalities=["Text", "Image"])
        )
        for part in response.candidates[0].content.parts:
            if part.inline_data is not None:
                image = Image.open(BytesIO(part.inline_data.data))
                image.save(os.path.join(OUTPUT_FOLDER, 'generated_chart.png'))
    except Exception as e:
        print("Image generation failed:", e)

# ---------- Markdown Utils ----------
def markdown_to_text(md_text):
    html = markdown2.markdown(md_text)
    soup = BeautifulSoup(html, "html.parser")
    return soup.get_text()

# ---------- PPTX Generator ----------
template_pptx = '/home/rootsh/Desktop/Mini Project/template_reports/Project status report.pptx'
def create_pptx_from_text(report_text, output_file_path, extracted_images, template_pptx):
    prs = Presentation(template_pptx) if template_pptx else Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    title_font = "Segoe UI"
    body_font = "Calibri Light"
    heading_color = RGBColor(0, 76, 153)
    body_color = RGBColor(40, 40, 40)

    def add_slide(title, lines):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
        tf = title_box.text_frame
        tf.text = title
        tf.paragraphs[0].font.size = Pt(36)
        tf.paragraphs[0].font.color.rgb = heading_color
        tf.paragraphs[0].font.name = title_font

        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(5))
        ctf = content_box.text_frame
        ctf.word_wrap = True

        for i, line in enumerate(lines):
            if i >= 12:
                break
            if line.strip():
                para = ctf.add_paragraph()
                para.text = line.strip()
                para.font.size = Pt(20)
                para.font.name = body_font
                para.font.color.rgb = body_color

    plain_text = markdown_to_text(report_text)
    sections = plain_text.split("\n\n")

    for sec in sections:
        lines = sec.strip().split("\n")
        if not lines:
            continue
        title = lines[0].strip()[:80]
        content_lines = lines[1:] if len(lines) > 1 else ["Content not available."]
        add_slide(title, content_lines)

    chart_path = os.path.join(OUTPUT_FOLDER, "generated_chart.png")
    if os.path.exists(chart_path):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tf = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1)).text_frame
        tf.text = "📊 Data Visualization"
        tf.paragraphs[0].font.size = Pt(36)
        tf.paragraphs[0].font.color.rgb = heading_color
        tf.paragraphs[0].font.name = title_font
        slide.shapes.add_picture(chart_path, Inches(2), Inches(1.5), width=Inches(9))

    for img_path in extracted_images:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tf = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1)).text_frame
        tf.text = "📸 Extracted Image"
        tf.paragraphs[0].font.size = Pt(36)
        tf.paragraphs[0].font.color.rgb = heading_color
        tf.paragraphs[0].font.name = title_font
        slide.shapes.add_picture(img_path, Inches(2), Inches(1.5), width=Inches(9))

    prs.save(output_file_path)

# ---------- Flask Routes ----------
@app.route("/generate-report", methods=["POST"])
def generate_report():
    if 'files' not in request.files:
        return jsonify({"error": "No files uploaded"}), 400

    files = request.files.getlist("files")
    template_reports = request.files.get('template')

    for folder in [UPLOAD_FOLDER, TEMPLATE_FOLDER]:
        for f in os.listdir(folder):
            os.remove(os.path.join(folder, f))

    extracted_images = []
    for file in files:
        file_path = os.path.join(UPLOAD_FOLDER, secure_filename(file.filename))
        file.save(file_path)

        if file.filename.endswith(".pdf"):
            extracted_images.extend(extract_images_from_pdf(file_path, OUTPUT_FOLDER))
        elif file.filename.endswith(".docx"):
            extracted_images.extend(extract_images_from_docx(file_path, OUTPUT_FOLDER))

    template_pptx = None
    headings = None
    if template_reports:
        filename = secure_filename(template_reports.filename)
        saved_path = os.path.join(TEMPLATE_FOLDER, filename)
        template_reports.save(saved_path)

        if filename.endswith(".pptx"):
            template_pptx = saved_path
        elif filename.endswith(".docx"):
            headings = extract_headings_from_docx(saved_path)

    combined_text = "\n".join([
        extract_text_from_pdf(os.path.join(UPLOAD_FOLDER, f)) if f.endswith(".pdf")
        else extract_text_from_docx(os.path.join(UPLOAD_FOLDER, f)) if f.endswith(".docx")
        else extract_text_from_excel(os.path.join(UPLOAD_FOLDER, f))
        for f in os.listdir(UPLOAD_FOLDER)
    ])

    report_text = generate_report_text_with_gemini(combined_text, headings)

    summary_prompt = f"Summarize the document contents in one paragraph for data visualization:\n\n{report_text[:6000]}"
    summary_response = client.models.generate_content(model='gemini-2.0-flash', contents=[summary_prompt])
    summary_text = summary_response.text.strip()

    generate_image_with_gemini(f"Generate a data visualization image based on this summary:\n{summary_text}")

    pptx_path = os.path.join(OUTPUT_FOLDER, "final_report.pptx")
    create_pptx_from_text(report_text, pptx_path, extracted_images, template_pptx)

    return jsonify({
        "download_url": "http://localhost:5000/download/final_report.pptx",
        "chart_image_url": "http://localhost:5000/download/generated_chart.png"
    })

@app.route("/download/<filename>")
def download_file(filename):
    return send_from_directory(OUTPUT_FOLDER, filename, as_attachment=True)

# ---------- Run ----------
if __name__ == "__main__":
    app.run(debug=True, port=5000)